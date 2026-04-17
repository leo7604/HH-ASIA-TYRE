"""
Create PowerPoint presentation for HH Asia Tyre Hackathon
Generates a .pptx file with all 12 slides
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Brand colors
    YELLOW = RGBColor(0xFF, 0xD7, 0x00)
    BLACK = RGBColor(0x00, 0x00, 0x00)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    GRAY = RGBColor(0x66, 0x66, 0x66)
    
    # Helper function to add slide with title and content
    def add_slide(title, content_lines, screenshot_placeholder=False):
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = BLACK
        p.alignment = PP_ALIGN.CENTER
        
        # Add content
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(6), Inches(5.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for i, line in enumerate(content_lines):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            p.text = line
            p.font.size = Pt(20)
            p.font.color.rgb = BLACK
            p.space_after = Pt(8)
        
        # Add screenshot placeholder
        if screenshot_placeholder:
            placeholder = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(7),
                Inches(1.5),
                Inches(5.5),
                Inches(4.5)
            )
            placeholder.fill.solid()
            placeholder.fill.fore_color.rgb = RGBColor(0xF0, 0xF0, 0xF0)
            placeholder.line.color.rgb = GRAY
            
            # Add text to placeholder
            text_box = slide.shapes.add_textbox(Inches(7), Inches(3.5), Inches(5.5), Inches(1))
            tf = text_box.text_frame
            p = tf.paragraphs[0]
            p.text = "📸 Insert Screenshot Here"
            p.font.size = Pt(18)
            p.font.color.rgb = GRAY
            p.alignment = PP_ALIGN.CENTER
    
    # SLIDE 1: Title
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "HH ASIA TYRE ALLIANCE PLUS+"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = BLACK
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.3), Inches(12.333), Inches(0.8))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Digital Booking System"
    p.font.size = Pt(32)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER
    
    # Team
    team_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12.333), Inches(1.5))
    tf = team_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Team: Leonardo Reparejo Jr., Asher Dane Manzano, Eric Godwin Jose, Felix Gallardo"
    p.font.size = Pt(20)
    p.font.color.rgb = BLACK
    p.alignment = PP_ALIGN.CENTER
    
    # School/Date placeholder
    info_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(12.333), Inches(0.8))
    tf = info_box.text_frame
    p = tf.paragraphs[0]
    p.text = "[School Name] | [Date]"
    p.font.size = Pt(18)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER
    
    # SLIDE 2: Introduction
    add_slide("ABOUT THE PROJECT", [
        "What We Built:",
        "✓ Modern landing page with promotions",
        "✓ 6-step booking system",
        "✓ Real-time slot availability",
        "✓ Multi-branch support (6 locations)",
        "✓ Mobile-responsive PWA",
        "",
        "Tech Stack:",
        "• React 18 + Vite 5",
        "• Supabase (PostgreSQL)",
        "• Tailwind CSS",
        "• Vercel Deployment",
        "• AI Assistant: Qwen (Qoder)",
        "",
        "Scope: 6 branches | 6+ services | ~8,000+ lines"
    ], True)
    
    # SLIDE 3: The Problem
    add_slide("THE CHALLENGE", [
        "Customer Pain Points:",
        "❌ Long, confusing booking forms (11 fields)",
        "❌ No real-time availability checking",
        "❌ Double bookings causing conflicts",
        "❌ Poor mobile experience",
        "❌ Unclear branch information",
        "",
        "Business Impact:",
        "❌ High booking abandonment rate",
        "❌ Frustrated customers",
        "❌ Lost revenue opportunities",
        "❌ Inefficient manual phone booking",
        "",
        "The Need: Modern, user-friendly online booking system"
    ], True)
    
    # SLIDE 4: Our Solution
    add_slide("THE SOLUTION", [
        "1. MODERN LANDING PAGE",
        "   • Hero slider with promotions",
        "   • Services showcase",
        "   • Branch locations",
        "",
        "2. GUIDED BOOKING FLOW",
        "   • 6 simple steps",
        "   • Real-time availability",
        "   • Instant confirmation",
        "",
        "3. CRM & INVENTORY INTEGRATION",
        "   • Customer data sync",
        "   • Real-time stock management",
        "   • Product availability",
        "",
        "4. MOBILE-FIRST DESIGN",
        "   • Progressive Web App",
        "   • Works on all devices"
    ], True)
    
    # SLIDE 5: Key Features
    add_slide("KEY FEATURES", [
        "1. MODERN LANDING PAGE",
        "   • Hero slider, services, branches",
        "",
        "2. SMART BOOKING SYSTEM",
        "   • 6-step process, real-time slots",
        "   • 1 customer per slot",
        "",
        "3. CRM & INVENTORY INTEGRATION",
        "   • Customer data sync",
        "   • Vehicle history tracking",
        "",
        "4. MOBILE-FIRST DESIGN",
        "   • PWA, touch-friendly",
        "",
        "Innovation: AI-assisted, Supabase + CRM"
    ], True)
    
    # SLIDE 6: Landing Page - Hero
    add_slide("LANDING PAGE: HERO SECTION", [
        "Features:",
        "✓ Eye-catching slider with promotions",
        "✓ Clear 'Book Now' call-to-action",
        "✓ Contact information visible",
        "✓ Professional branding",
        "",
        "Screenshot:",
        "• Full homepage hero section",
        "• Show promotions slider",
        "• Highlight 'Book Now' button"
    ], True)
    
    # SLIDE 7: Landing Page - Services
    add_slide("LANDING PAGE: SERVICES & BRANCHES", [
        "Services Section:",
        "✓ 6+ automotive services",
        "✓ Service icons and descriptions",
        "✓ Price range estimates",
        "",
        "Branch Locations:",
        "✓ All 6 branches with images",
        "✓ Region filtering (Manila vs Laoag)",
        "✓ Contact details and hours",
        "",
        "Screenshot:",
        "• Services grid view",
        "• Branch cards with images"
    ], True)
    
    # SLIDE 8: Landing Page - Testimonials
    add_slide("LANDING PAGE: SOCIAL PROOF", [
        "Testimonials:",
        "✓ Customer reviews",
        "✓ Star ratings",
        "✓ Trust signals",
        "",
        "Gallery:",
        "✓ Service photos",
        "✓ Workshop images",
        "✓ Professional quality",
        "",
        "Screenshot:",
        "• Testimonials section",
        "• Gallery grid"
    ], True)
    
    # SLIDE 9: How to Book - Part 1
    add_slide("HOW TO BOOK: Getting Started", [
        "Step 1: Select Branch",
        "• Choose region (Manila/Laoag)",
        "• View branch details",
        "",
        "Step 2: Enter Vehicle",
        "• Year, make, model",
        "• Plate number",
        "",
        "Step 3: Choose Services",
        "• Multi-select services",
        "• See descriptions and prices",
        "",
        "Screenshots:",
        "• Branch selection, Vehicle form,",
        "• Service checkboxes"
    ], True)
    
    # SLIDE 10: How to Book - Part 2
    add_slide("HOW TO BOOK: Complete Booking", [
        "Step 4: Pick Date & Time",
        "• Calendar date picker",
        "• Color-coded availability",
        "  🟢 Available | 🟡 Limited |  Full",
        "",
        "Step 5: Customer Information",
        "• Name and contact",
        "• Email and phone",
        "",
        "Step 6: Review & Confirm",
        "• Complete summary",
        "• Instant reference number",
        "",
        "Screenshot:",
        "• Date/time picker, Info form,",
        "• Success page"
    ], True)
    
    # SLIDE 11: Mobile Experience
    add_slide("MOBILE-FIRST DESIGN", [
        "Progressive Web App (PWA):",
        "✓ Installable on mobile",
        "✓ Works offline",
        "✓ App-like experience",
        "",
        "Responsive Design:",
        "✓ Desktop optimized",
        "✓ Tablet friendly",
        "✓ Mobile optimized",
        "• Touch-friendly buttons",
        "",
        "Screenshots:",
        "• Homepage (mobile)",
        "• Booking form (mobile)",
        "• Time slots (mobile)"
    ], True)
    
    # SLIDE 12: Conclusion
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "CONCLUSION"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = BLACK
    p.alignment = PP_ALIGN.CENTER
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12.333), Inches(3))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    lines = [
        "Achievements:",
        "✅ Modern landing page",
        "✅ Complete booking system",
        "✅ Real-time availability",
        "✅ Multi-branch support",
        "✅ Mobile-responsive PWA",
        "",
        "Impact:",
        "• Improved customer experience",
        "• Efficient online booking",
        "• Professional web presence",
        "",
        "Thank You! Questions?"
    ]
    
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        p.text = line
        p.font.size = Pt(28)
        p.font.color.rgb = BLACK
        p.space_after = Pt(10)
        p.alignment = PP_ALIGN.CENTER
    
    # Save presentation
    prs.save('HH_Asia_Tyre_Presentation.pptx')
    print("✅ PowerPoint created: HH_Asia_Tyre_Presentation.pptx")
    print("📊 Total slides: 12")
    print("⏱️ Estimated time: 12 minutes")
    print("\n📸 Remember to add screenshots to each slide!")

if __name__ == "__main__":
    create_presentation()
